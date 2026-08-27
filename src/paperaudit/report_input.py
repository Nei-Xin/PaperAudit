"""Parsing adapters for user-provided report files."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Literal
from zipfile import BadZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError


_MAX_PPTX_BYTES = 100 * 1024 * 1024
_MAX_SLIDES = 300


@dataclass(frozen=True)
class ParsedReportInput:
    text: str
    kind: Literal["text", "pptx"]
    page_count: int | None = None
    warnings: tuple[str, ...] = ()


def decode_report_bytes(data: bytes) -> str:
    if not data:
        raise ValueError("报告文件为空。")
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            text = data.decode(encoding)
        except UnicodeDecodeError:
            continue
        if text.strip():
            return text
        raise ValueError("报告文件中没有可审计的文本。")
    raise ValueError("报告文件编码无法识别，请保存为 UTF-8 后重试。")


def parse_report_file(data: bytes, filename: str) -> ParsedReportInput:
    """Parse one supported report file without changing its source bytes."""

    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md"}:
        return ParsedReportInput(text=decode_report_bytes(data), kind="text")
    if suffix == ".pptx":
        return parse_pptx_report(data)
    raise ValueError("暂不支持该报告格式，请上传 .txt、.md 或 .pptx 文件。")


def parse_pptx_report(data: bytes) -> ParsedReportInput:
    if not data:
        raise ValueError("PowerPoint 文件为空。")
    if len(data) > _MAX_PPTX_BYTES:
        raise ValueError("PowerPoint 文件超过 100MB，请压缩图片后重试。")
    try:
        presentation = Presentation(BytesIO(data))
    except (PackageNotFoundError, BadZipFile, ValueError, KeyError) as exc:
        raise ValueError("PowerPoint 文件损坏或不是有效的 .pptx 文件。") from exc

    slide_count = len(presentation.slides)
    if not slide_count:
        raise ValueError("PowerPoint 中没有幻灯片。")
    if slide_count > _MAX_SLIDES:
        raise ValueError("PowerPoint 超过 300 页，请拆分后再审计。")

    sections: list[str] = ["# PowerPoint 审计输入"]
    visual_slides: list[int] = []
    content_slide_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = _shape_text(title_shape) if title_shape is not None else ""
        body_lines: list[str] = []
        table_lines: list[str] = []
        chart_lines: list[str] = []
        picture_count = 0

        shapes = sorted(slide.shapes, key=lambda shape: (shape.top, shape.left))
        for shape in shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            shape_body, shape_tables, shape_charts, shape_pictures = _extract_shape(shape)
            body_lines.extend(shape_body)
            table_lines.extend(shape_tables)
            chart_lines.extend(shape_charts)
            picture_count += shape_pictures

        notes = _slide_notes(slide)
        if picture_count:
            visual_slides.append(slide_number)
        if title or body_lines or table_lines or chart_lines or notes:
            content_slide_count += 1

        slide_parts = [f"[幻灯片 {slide_number}]"]
        if title:
            slide_parts.append(f"标题：{title}")
        if body_lines:
            slide_parts.extend(["正文：", *body_lines])
        if table_lines:
            slide_parts.extend(["表格：", *table_lines])
        if chart_lines:
            slide_parts.extend(["图表数据：", *chart_lines])
        if notes:
            slide_parts.extend(["演讲者备注：", notes])
        sections.append("\n".join(slide_parts))

    if not content_slide_count:
        raise ValueError(
            "PowerPoint 中没有可提取的文本或表格；当前版本尚不能单独审计纯图片幻灯片。"
        )

    warnings: list[str] = []
    if visual_slides:
        pages = "、".join(str(page) for page in visual_slides[:12])
        suffix = "等" if len(visual_slides) > 12 else ""
        warnings.append(
            f"第 {pages} 页{suffix}包含图片；当前先审计可提取文本，图片内容需视觉能力复核。"
        )
    return ParsedReportInput(
        text="\n\n".join(sections).strip(),
        kind="pptx",
        page_count=slide_count,
        warnings=tuple(warnings),
    )


def _extract_shape(shape: object) -> tuple[list[str], list[str], list[str], int]:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        body: list[str] = []
        tables: list[str] = []
        charts: list[str] = []
        pictures = 0
        for child in sorted(shape.shapes, key=lambda item: (item.top, item.left)):
            child_body, child_tables, child_charts, child_pictures = _extract_shape(child)
            body.extend(child_body)
            tables.extend(child_tables)
            charts.extend(child_charts)
            pictures += child_pictures
        return body, tables, charts, pictures

    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            values = [_normalize_text(cell.text) for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        return [], rows, [], 0

    if getattr(shape, "has_chart", False):
        return [], [], _chart_text(shape.chart), 0

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return [], [], [], 1

    text = _shape_text(shape)
    return ([f"- {text}"] if text else []), [], [], 0


def _shape_text(shape: object) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = _normalize_text(paragraph.text)
        if text:
            indent = "  " * min(getattr(paragraph, "level", 0), 3)
            paragraphs.append(f"{indent}{text}")
    return "\n".join(paragraphs)


def _chart_text(chart: object) -> list[str]:
    lines: list[str] = []
    try:
        if chart.has_title:
            title = _normalize_text(chart.chart_title.text_frame.text)
            if title:
                lines.append(f"图表标题：{title}")
        for series in chart.series:
            name = _normalize_text(str(series.name)) or "未命名系列"
            values = ["" if value is None else str(value) for value in series.values]
            lines.append(f"系列 {name}：{', '.join(values)}")
    except (AttributeError, TypeError, ValueError):
        return lines
    return lines


def _slide_notes(slide: object) -> str:
    try:
        text_frame = slide.notes_slide.notes_text_frame
    except (AttributeError, KeyError, ValueError):
        return ""
    return _normalize_text(text_frame.text) if text_frame is not None else ""


def _normalize_text(text: str) -> str:
    return " ".join(text.replace("\x0b", "\n").split())
