from __future__ import annotations

from io import BytesIO

import pytest
from pptx import Presentation
from pptx.util import Inches

from paperaudit.report_input import decode_report_bytes, parse_report_file


def test_decode_report_bytes_supports_utf8_bom_and_gb18030() -> None:
    assert decode_report_bytes("中文报告".encode("utf-8-sig")) == "中文报告"
    assert decode_report_bytes("审计内容".encode("gb18030")) == "审计内容"


@pytest.mark.parametrize("payload", [b"", b"   \n\t"])
def test_decode_report_bytes_rejects_empty_text(payload: bytes) -> None:
    with pytest.raises(ValueError, match="为空|没有可审计"):
        decode_report_bytes(payload)


def test_decode_report_bytes_rejects_unknown_encoding() -> None:
    with pytest.raises(ValueError, match="编码无法识别"):
        decode_report_bytes(b"\x81")


def _sample_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "实验结果"
    slide.placeholders[1].text = "SAM 在 23 个数据集中的 16 个上优于 RITM。"
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(4), Inches(5), Inches(1)
    ).table
    table.cell(0, 0).text = "方法"
    table.cell(0, 1).text = "IoU"
    table.cell(1, 0).text = "SAM"
    table.cell(1, 1).text = "81.2"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_parse_report_file_extracts_pptx_slide_text_and_table() -> None:
    payload = _sample_pptx()

    parsed = parse_report_file(payload, "report.PPTX")

    assert parsed.kind == "pptx"
    assert parsed.page_count == 1
    assert "[幻灯片 1]" in parsed.text
    assert "标题：实验结果" in parsed.text
    assert "SAM 在 23 个数据集中的 16 个上优于 RITM" in parsed.text
    assert "方法 | IoU" in parsed.text
    assert "SAM | 81.2" in parsed.text


def test_parse_report_file_keeps_plain_text_compatibility() -> None:
    parsed = parse_report_file("中文报告".encode(), "report.md")

    assert parsed.kind == "text"
    assert parsed.page_count is None
    assert parsed.text == "中文报告"


def test_parse_report_file_rejects_invalid_pptx() -> None:
    with pytest.raises(ValueError, match="损坏|有效"):
        parse_report_file(b"not a pptx", "report.pptx")


def test_parse_report_file_rejects_unknown_extension() -> None:
    with pytest.raises(ValueError, match="暂不支持"):
        parse_report_file(b"data", "report.docx")
